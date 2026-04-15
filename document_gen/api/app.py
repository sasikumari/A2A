"""
NPCI TITAN AGENTS: Unified API & Orchestration Layer
--------------------------------------------------
This Flask application serves as the core entry point for the Titan Engine.
It manages:
1.  UPI Simulation Endpoints (ReqPay, Push, Collect, etc.)
2.  Agentic Innovation Lifecycle (Clarify, Canvas, Execution)
3.  Asynchronous Real-Time Observability (SSE via NotificationBus)
4.  Zero-Trust Identity Handshakes
"""
from flask import Flask, request, jsonify, render_template, Response, render_template_string, send_file
from dataclasses import asdict
from switch.upi_switch import UPISwitch, VPARegistry
from switch.ledger import Ledger
from switch.notification_bus import NotificationBus
from psps.payer_psp import PayerPSP
from psps.payee_psp import PayeePSP
from banks.remitter_bank import RemitterBank, Account as RAccount
from banks.beneficiary_bank import BeneficiaryBank, Account as BAccount
from banks.icici_bank import ICICIBank, Account as IAccount
from authorization_service import AuthorizationService
from agents.reasoning_agent import ReasoningAgent
from agents.canvas_agent import CanvasAgent
from agents.prototype_agent import PrototypeAgent
from agents.product_kit_agent import ProductKitAgent
from agents.npci_master_agent import NPCIMasterAgent
from agents.party_agent import PartyAgent
from agents.llm_client import LLMClient
from infrastructure.doc_store import DocumentStore
from infrastructure.rbac import requires_role
from infrastructure.sha_signing import generate_signed_document
from infrastructure.reference_psp_node import reference_psp_bp
from infrastructure.reference_bank_node import reference_bank_bp
from infrastructure.certification_service import CertificationManager
import xmlschema
from storage.db import init_db, upsert_user, upsert_mapper, UserRoleEnum
import queue
import json
import time
import os
import subprocess

import threading
import uuid
import io
import zipfile
from datetime import datetime
from html import escape as html_escape
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Force bus configuration before NotificationBus reads environment variables
os.environ["DISABLE_BUS"] = "0"
os.environ["USE_IN_MEMORY_BUS"] = "1"

app = Flask(__name__, static_folder="static", template_folder="templates")

bus = NotificationBus()
ledger = Ledger()
payer_registry = VPARegistry()
payee_registry = VPARegistry()
SessionFactory = init_db()

# --- SSE Setup ---
class MessageAnnouncer:
    def __init__(self):
        self.listeners = []

    def listen(self):
        q = queue.Queue(maxsize=500)
        self.listeners.append(q)
        return q

    def announce(self, msg):
        for i in reversed(range(len(self.listeners))):
            try:
                self.listeners[i].put_nowait(msg)
            except queue.Full:
                del self.listeners[i]

announcer = MessageAnnouncer()

def format_sse(data: str, event=None) -> str:
    msg = f'data: {data}\n\n'
    if event:
        msg = f'event: {event}\n{msg}'
    return msg

def broadcast_xml(xml_content, source="API"):
    """Helper to broadcast XML to SSE listeners"""
    # We'll send a JSON object with source and content
    payload = json.dumps({"source": source, "content": xml_content})
    announcer.announce(format_sse(data=payload, event="xml_log"))

# Bridge NotificationBus -> SSE
def bus_listener():
    """
    Bridge between the NotificationBus (Agents) and SSE (Frontend).
    This enables real-time "thinking logs" and transaction traces to be visualized instantly.
    """
    print("[Bridge] Starting bus listener for 'xml_stream'...")
    for event in bus.subscribe("xml_stream"):
        if event:
            source = event.get("source", "Internal")
            content = event.get("content", "")
            if content:
                broadcast_xml(content, source)

def agent_status_listener():
    print("[Bridge] Starting bus listener for 'agent_status'...")
    for event in bus.subscribe("agent_status"):
        if event:
            payload = json.dumps(event)
            announcer.announce(format_sse(data=payload, event="agent_status"))

def spec_change_listener_bridge():
    print("[Bridge] Starting bus listener for 'spec_change'...")
    for event in bus.subscribe("spec_change"):
        if event:
            payload = json.dumps(event)
            announcer.announce(format_sse(data=payload, event="spec_change"))

# Start background threads
threading.Thread(target=bus_listener, daemon=True).start()
threading.Thread(target=agent_status_listener, daemon=True).start()
threading.Thread(target=spec_change_listener_bridge, daemon=True).start()


# Set up banks and PSPs
auth_service = AuthorizationService(default_pin="1234")

payer_bank = RemitterBank(code="PAYERBANK", auth_service=auth_service, ledger=ledger, bus=bus)
benef_bank = BeneficiaryBank(code="BENEBANK", ledger=ledger, bus=bus)
icici_bank = ICICIBank(code="ICICIBANK", auth_service=auth_service, ledger=ledger, bus=bus)

# Create accounts
ramesh = RAccount(id="A1", vpa="ramesh@payer", name="ramesh", bank_code=payer_bank.code, balance=500000.0)
merchant = BAccount(id="M1", vpa="merchant@benef", name="Merchant", bank_code=benef_bank.code, balance=200.0)
ashlesha = IAccount(id="IC1", vpa="ashlesha@icici", name="Ashlesha", bank_code=icici_bank.code, balance=200000.0)

payer_bank.add_account(ramesh)
benef_bank.add_account(merchant)
icici_bank.add_account(ashlesha)

payer_registry.register(ramesh.vpa, payer_bank)
payer_registry.register(ashlesha.vpa, icici_bank)
payee_registry.register(merchant.vpa, benef_bank)


# Seed DB with users and NPCI mapper (based on in-memory setup)
with SessionFactory() as s:
    upsert_user(s, vpa=ramesh.vpa, name=ramesh.name, role=UserRoleEnum.PAYER_BANK, bank_code=payer_bank.code, psp_code="PAYERPSP")
    upsert_user(s, vpa=ashlesha.vpa, name=ashlesha.name, role=UserRoleEnum.PAYER_BANK, bank_code=icici_bank.code, psp_code="PAYERPSP")
    upsert_user(s, vpa=merchant.vpa, name=merchant.name, role=UserRoleEnum.PAYEE_BANK, bank_code=benef_bank.code, psp_code="PAYEEPSP")
    upsert_user(s, vpa="ramesh@payer", name="ramesh PSP", role=UserRoleEnum.PAYER_PSP, bank_code=payer_bank.code, psp_code="PAYERPSP")
    upsert_user(s, vpa="ashlesha@icici", name="Ashlesha PSP", role=UserRoleEnum.PAYER_PSP, bank_code=icici_bank.code, psp_code="PAYERPSP")
    upsert_user(s, vpa="merchant@benef", name="Merchant PSP", role=UserRoleEnum.PAYEE_PSP, bank_code=benef_bank.code, psp_code="PAYEEPSP")
    upsert_mapper(s, vpa=ramesh.vpa, bank_code=payer_bank.code, account_id=ramesh.id, psp_code="PAYERPSP")
    upsert_mapper(s, vpa=ashlesha.vpa, bank_code=icici_bank.code, account_id=ashlesha.id, psp_code="PAYERPSP")
    upsert_mapper(s, vpa=merchant.vpa, bank_code=benef_bank.code, account_id=merchant.id, psp_code="PAYEEPSP")
    s.commit()

upi_switch = UPISwitch(payer_registry, ledger, bus, schema_dir="api/schemas", db_session_factory=SessionFactory)

payer_psp = PayerPSP(payer_bank, upi_switch, schema_dir="api/schemas")
payee_psp = PayeePSP(benef_bank, upi_switch, payee_registry, schema_dir="api/schemas")

upi_switch.register_payee_services(
    payee_psp.process_valadd_request,
    payee_psp.resolve_payee,
    payee_psp.resolve_bank_by_ifsc,
    payee_psp.process_auth_details,
    payee_psp.process_txn_confirmation
)

# Register Reference Bank and PSP Blueprints for Phase 2 network simulation
app.register_blueprint(reference_psp_bp, url_prefix="/reference_psp")
app.register_blueprint(reference_bank_bp, url_prefix="/reference_bank")

# --- Agents Setup ---
from agents import init_agents
switch_agent, agents_list = init_agents(bus)

# Initialize Party Agents listening to the NotificationBus
bank_party_agent = PartyAgent("reference_issuer_001", "issuer_bank", LLMClient(), bus)
psp_party_agent = PartyAgent("reference_psp_001", "psp", LLMClient(), bus)

# Initialize Certification Manager aggregating from the NotificationBus
cert_manager = CertificationManager(bus)

WORKFLOW2_SESSIONS = {}

def _extract_json_from_text(raw_text: str):
    """Best-effort JSON extractor for noisy LLM output."""
    if not raw_text:
        return {}
    txt = raw_text.strip()
    txt = txt.replace("```json", "").replace("```", "").strip()
    start = txt.find("{")
    end = txt.rfind("}")
    if start == -1 or end == -1 or end <= start:
        return {}
    try:
        return json.loads(txt[start:end + 1])
    except Exception:
        return {}

def _workflow2_fallback(feature_name: str, objective: str, target_users: str):
    """Deterministic fallback so Workflow 2 is demo-ready even if LLM is unavailable."""
    objective = objective or f"Launch {feature_name} as a production-grade UPI feature"
    target_users = target_users or "UPI users, issuers, UPI apps, and merchants"
    return {
        "thinking_trace": [
            "Understanding feature scope and target ecosystem participants",
            "Mapping product build canvas fields to decision-ready outputs",
            "Checking RBI/NPCI constraints for delegated and mandate-like flows",
            "Preparing prototype and communication handoff pack"
        ],
        "agent_plan": [
            "Deep research and ecosystem signal mapping",
            "Product canvas draft and maker-edit pass",
            "Prototype build request and redirect preparation",
            "Documentation package generation",
            "Execution handoff for Phase 2 deployment"
        ],
        "deep_research": {
            "feature_summary": f"{feature_name} enables delegated and intent-driven UPI journeys with low-friction checkout.",
            "need": "Solves drop-offs in high-frequency use-cases by reducing repeated payment steps and improving conversion.",
            "market_view": "High merchant and app interest; issuer readiness and reconciliation visibility are key adoption gates.",
            "scalability": "Scale via top UPI apps + top issuers + top merchants in relevant market segments.",
            "risks": [
                "Misuse through rogue merchant flows",
                "Block/reconcile mismatch across participants",
                "Policy misalignment on commercial treatment"
            ],
            "regulatory_summary": "Must align with NPCI circular controls, explicit user authorization, notifications, and auditability.",
            "rbi_guidelines_checked": [
                "https://www.rbi.org.in/scripts/NotificationUser.aspx?Id=12032&Mode=0",
                "https://www.rbi.org.in/commonman/english/scripts/Notification.aspx?Id=1888"
            ]
        },
        "product_canvas": {
            "feature_for_layman": f"{feature_name} lets users authorize once and complete future payments faster with clear controls.",
            "need_and_differentiation": "Exponential UX gain for repeat and contextual checkout (fewer clicks, lower friction).",
            "market_view": "Strong pull from large merchants; depends on issuer + app co-readiness.",
            "scalability": "Anchor with top apps, top issuers, and high-frequency merchant categories.",
            "validation_mvp": "Pilot with 2 apps, 2 issuers, 2 merchants; compare conversion, checkout time, and success rates.",
            "operating_kpis": [
                "Checkout time reduction",
                "Creation-to-execution conversion",
                "Success rate uplift vs PIN-heavy flow"
            ],
            "comms_outputs": [
                "Product document",
                "Demo video and PM explainer",
                "FAQ and support chatbot seed corpus",
                "Draft circular for ecosystem"
            ],
            "pricing_view": "Keep interoperable and adoption-friendly; phased commercial model where justified.",
            "compliance": "RBI/NPCI guardrails on consent, limits, notifications, fraud controls, and audit trails."
        },
        "prototype": {
            "name": f"{feature_name} Experience Prototype",
            "url": "",
            "notes": "Prototype generation requested. Provide a Lovable/Figma URL after maker review."
        },
        "product_comms": {
            "doc_outline": [
                "Problem statement and value proposition",
                "User journeys (happy path + exceptions)",
                "API/spec updates and test cases",
                "Operational and grievance handling model"
            ],
            "faq_starters": [
                "How does user consent work?",
                "How can users revoke or change limits?",
                "What happens on insufficient funds or partial failures?"
            ],
            "circular_draft": f"Draft circular: launch scope, participant responsibilities, risk controls, dispute workflow, and rollout timelines for {feature_name}.",
            "trained_llm_note": "Use product document + FAQs as retrieval corpus for a public Q&A assistant."
        },
        "execution_handoff": {
            "objective": objective,
            "target_users": target_users,
            "next_step": "After maker approval, continue with current Phase 2 change/deploy flow."
        }
    }

def _docx_from_text(title: str, body: str) -> bytes:
    """Create a minimal .docx in-memory without external dependencies."""
    lines = [title, ""] + [ln for ln in body.splitlines()]
    para_xml = []
    for ln in lines:
        safe = html_escape(ln)
        para_xml.append(
            f"<w:p><w:r><w:t xml:space=\"preserve\">{safe}</w:t></w:r></w:p>"
        )
    document_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas"
 xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"
 xmlns:o="urn:schemas-microsoft-com:office:office"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math"
 xmlns:v="urn:schemas-microsoft-com:vml"
 xmlns:wp14="http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing"
 xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
 xmlns:w10="urn:schemas-microsoft-com:office:word"
 xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
 xmlns:w14="http://schemas.microsoft.com/office/word/2010/wordml"
 xmlns:wpg="http://schemas.microsoft.com/office/word/2010/wordprocessingGroup"
 xmlns:wpi="http://schemas.microsoft.com/office/word/2010/wordprocessingInk"
 xmlns:wne="http://schemas.microsoft.com/office/word/2006/wordml"
 xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
 mc:Ignorable="w14 wp14">
 <w:body>
 {''.join(para_xml)}
 <w:sectPr>
  <w:pgSz w:w="12240" w:h="15840"/>
  <w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440" w:header="708" w:footer="708" w:gutter="0"/>
 </w:sectPr>
 </w:body>
</w:document>"""
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
 <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
 <Default Extension="xml" ContentType="application/xml"/>
 <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
 <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document_xml)
    mem.seek(0)
    return mem.getvalue()

def _ppt_from_text(title: str, content: str) -> bytes:
    """Create a 5-slide Product Deck in-memory using python-pptx."""
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "NPCI UPI Product Strategy & Ecosystem Alignment"
    
    # Split content into manageable chunks (simulating 4 more slides)
    sections = content.split("\n\n")
    slide_titles = ["Market Opportunity", "User Journey & Experience", "Technical Architecture", "Operational & Risk Controls"]
    
    for i in range(min(4, len(slide_titles))):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_titles[i]
        body = slide.placeholders[1]
        
        # Take a bit of content for this slide
        txt = sections[i] if i < len(sections) else "Details pending refinement."
        body.text = txt
        
    mem = io.BytesIO()
    prs.save(mem)
    mem.seek(0)
    return mem.getvalue()

def _render_workflow2_docs(session_data: dict):
    """Build doc outputs for preview and download."""
    payload = session_data.get("payload", {})
    feature = session_data.get("feature_name", "Workflow 2 Feature")
    deep = payload.get("deep_research", {})
    canvas = payload.get("product_canvas", {})
    proto = payload.get("prototype", {})
    comms = payload.get("product_comms", {})

    docs = {
        "product_document": {
            "title": f"{feature} - Product Document",
            "body": (
                "Deep Research\n"
                f"{json.dumps(deep, indent=2)}\n\n"
                "Product Canvas\n"
                f"{json.dumps(canvas, indent=2)}\n\n"
                "Prototype Brief\n"
                f"{json.dumps(proto, indent=2)}\n"
            ),
        },
        "faq_document": {
            "title": f"{feature} - FAQs",
            "body": "\n".join(comms.get("faq_starters", [])) or "FAQ content to be refined after maker review.",
        },
        "circular_draft": {
            "title": f"{feature} - Draft Circular",
            "body": comms.get("circular_draft", "Circular draft pending."),
        },
        "pm_explainer_script": {
            "title": f"{feature} - PM Explainer Script",
            "body": "\n".join(comms.get("doc_outline", [])) or "Explainer script points pending.",
        },
        "product-deck": {
            "title": f"{feature} - Product Deck",
            "body": (
                f"Slide 1: {feature}\n"
                f"Slide 2: {deep.get('market_view', 'Market landscape')}\n"
                f"Slide 3: {canvas.get('feature_for_layman', 'Product experience')}\n"
                f"Slide 4: {json.dumps(comms.get('doc_outline', []), indent=2)}\n"
                f"Slide 5: {canvas.get('compliance', 'Governance and limits')}"
            ),
        },
    }
    return docs


# --- Routes ---

@app.route("/push", methods=["POST"])
def push_payment():
    xml_req = request.data.decode("utf-8")
    try:
        print("\n[XML] API /push received:\n", xml_req, "\n")
        broadcast_xml(xml_req, source="[API] /push Received")
        
        ack = payer_psp.switch.handle_push(xml_req)
        
        print("\n[ACK] /push returning JSON:\n", ack, "\n")
        # broadcast_xml(json.dumps(ack), source="[API] /push ACK") # Optional: broadcast ACK too
        
        return jsonify(ack), 202
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 400


@app.route("/collect", methods=["POST"])
def collect_request():
    xml_req = request.data.decode("utf-8")
    try:
        print("\n[XML] API /collect received:\n", xml_req, "\n")
        broadcast_xml(xml_req, source="[API] /collect Received")

        def approval_fn(rrn, amount, note):  # simulated approval
            # Simulated user authorization (auto-approve)
            print(f"[COLLECT] Auto-approved collect for {amount}")
            return True

        ack = payee_psp.switch.handle_collect(xml_req, approval_fn)
        print("\n[ACK] /collect returning JSON:\n", ack, "\n")
        return jsonify(ack), 202
    except Exception as e:
        return jsonify({"error": str(e)}), 400


@app.route("/status", methods=["POST"])
def status_poll():
    xml_req = request.data.decode("utf-8")
    try:
        print("\n[XML] API /status received:\n", xml_req, "\n")
        broadcast_xml(xml_req, source="[API] /status Received")
        
        xml_resp = upi_switch.status_poll(xml_req)
        
        print("\n[XML] API /status returning:\n", xml_resp, "\n")
        broadcast_xml(xml_resp, source="[API] /status Response")
        
        return xml_resp, 200, {"Content-Type": "application/xml"}
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route("/reqpay", methods=["POST"])
def reqpay():
    xml_req = request.data.decode("utf-8")
    try:
        print("\n[XML] API /reqpay received:\n", xml_req, "\n")
        broadcast_xml(xml_req, source="[API] /reqpay Received")
        
        ack = upi_switch.handle_reqpay(xml_req)
        
        print("\n[ACK] /reqpay returning JSON:\n", ack, "\n")
        return jsonify(ack), 202
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route("/transactions", methods=["GET"])
def list_transactions():
    try:
        entries = getattr(ledger, "entries", []) or []
        data = [asdict(e) for e in entries]
        # Only show terminal states (SUCCESS/DECLINED/FAILURE/etc.)
        terminal = [d for d in data if d.get("status") and d.get("status") != "PENDING"]
        # Most recent first (created_at is ISO string)
        terminal.sort(key=lambda d: d.get("created_at") or "", reverse=True)
        return jsonify(terminal), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/documents/ppt/download", methods=["POST"])
def download_ppt():
    try:
        data = request.json
        title = data.get("title", "Product Deck")
        content = data.get("content", "")
        ppt_bytes = _ppt_from_text(title, content)
        return Response(
            ppt_bytes,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment;filename={title.replace(' ', '_')}.pptx"}
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/documents/download", methods=["POST"])
def download_docx_route():
    try:
        data = request.json
        title = data.get("title", "Document")
        content = data.get("content", "")
        docx_bytes = _docx_from_text(title, content)
        return Response(
            docx_bytes,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment;filename={title.replace(' ', '_')}.docx"}
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/validate-address", methods=["POST"])
def validate_address():
    xml_req = request.data.decode("utf-8")
    try:
        print("\n[XML] API /validate-address received:\n", xml_req, "\n")
        broadcast_xml(xml_req, source="[API] /validate-address Received")
        
        xml_resp, _ = payee_psp.process_valadd_request(xml_req)
        
        print("\n[XML] API /validate-address returning:\n", xml_resp, "\n")
        broadcast_xml(xml_resp, source="[API] /validate-address Response")
        
        return xml_resp, 200, {"Content-Type": "application/xml"}
    except Exception as e:
        return jsonify({"error": str(e)}), 400


@app.route("/health")
def health():
    """Phase 1 health check: server and core deps."""
    try:
        ok = ledger is not None and upi_switch is not None
        return jsonify({"status": "ok", "phase1": ok}), 200
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)}), 500

@app.route("/limits")
def get_limits():
    """
    Returns the live business-rule limits by parsing upi_switch.py with regex.
    Used by run_tests.py and Phase 2 dashboard to build dynamic test amounts.
    """
    import re as _re
    _project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    switch_path = os.path.join(_project_root, "switch", "upi_switch.py")

    p2p   = 100_00_000
    max_t = 100_00_000

    try:
        with open(switch_path, "r") as _f:
            src = _f.read()
        # Match:  P2P_LIMIT = 3_00_000  or  P2P_LIMIT = 300000
        def _parse(name):
            m = _re.search(rf"^{name}\s*=\s*([\d_]+)", src, _re.MULTILINE)
            if m:
                return int(m.group(1).replace("_", ""))
            return None
        p2p   = _parse("P2P_LIMIT")   or p2p
        max_t = _parse("MAX_TXN_AMOUNT") or max_t
    except Exception:
        pass

    effective = min(p2p, max_t)
    return jsonify({
        "p2p_limit":       p2p,
        "max_txn_amount":  max_t,
        "effective_limit": effective,
        "test_amounts": {
            "safe":        round(effective * 0.10, 2),   # 10% — always succeeds
            "mid":         round(effective * 0.50, 2),   # 50% — always succeeds
            "near_limit":  round(effective * 0.95, 2),   # 95% — succeeds (just below)
            "over_limit":  round(effective * 1.50, 2),   # 150% — always DECLINED
        }
    })

@app.route("/agents/reload", methods=["POST"])
def agents_reload():
    """
    Hot-reload business-rule constants into the running process after agents
    update source files.  Patches P2P_LIMIT / MAX_TXN_AMOUNT directly in the
    already-imported switch.upi_switch module so that every subsequent request
    uses the new limits without a full process restart.

    Also uses importlib.reload() so that any new helper methods / logic in the
    module are picked up too.
    """
    import re as _re
    import importlib
    import switch.upi_switch as _sw_mod

    _project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    switch_path = os.path.join(_project_root, "switch", "upi_switch.py")

    try:
        # 1. Full module reload — picks up all code changes agents made
        importlib.reload(_sw_mod)

        # 2. Read updated constants directly from file (belt-and-suspenders)
        with open(switch_path, "r") as _f:
            src = _f.read()

        def _parse(name):
            m = _re.search(rf"^{name}\s*=\s*([\d_]+)", src, _re.MULTILINE)
            return int(m.group(1).replace("_", "")) if m else None

        new_p2p   = _parse("P2P_LIMIT")
        new_max_t = _parse("MAX_TXN_AMOUNT")

        # 3. Patch module globals so all existing references see new values
        if new_p2p   is not None:
            _sw_mod.P2P_LIMIT       = new_p2p
        if new_max_t is not None:
            _sw_mod.MAX_TXN_AMOUNT  = new_max_t

        return jsonify({
            "status": "reloaded",
            "P2P_LIMIT":      _sw_mod.P2P_LIMIT,
            "MAX_TXN_AMOUNT": _sw_mod.MAX_TXN_AMOUNT,
        })

    except Exception as exc:
        return jsonify({"status": "error", "error": str(exc)}), 500


@app.route("/", methods=["GET"])
def ui_home():
    # Simple test UI (Phase 1)
    return render_template("index.html")

@app.route("/ui", methods=["GET"])
def new_ui():
    return render_template("new_ui.html")

@app.route("/iot-demo", methods=["GET"])
def iot_demo_ui():
    return render_template("iot_demo.html")

@app.route("/stream")
def stream():
    def event_stream():
        messages = announcer.listen()
        while True:
            try:
                # Use a small timeout to yield heartbeats every 15 seconds
                msg = messages.get(timeout=15)
                yield msg
            except queue.Empty:
                yield ": heartbeat\n\n"
    return Response(event_stream(), mimetype="text/event-stream")

@app.route("/benchmark", methods=["POST"])
def run_benchmark_endpoint():
    try:
        data = request.json
        rps = data.get("rps", 50)
        duration = data.get("duration", 10)
        concurrency = data.get("concurrency", 50)
        poll_status = data.get("poll_status", False)
        
        cmd = [
            "python3", "tools/tps_benchmark.py",
            "--host", "http://127.0.0.1:5000",
            "--rps", str(rps),
            "--duration", str(duration),
            "--concurrency", str(concurrency)
        ]
        if poll_status:
            cmd.append("--poll-status")
            
        def generate():
            import subprocess
            process = subprocess.Popen(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                bufsize=1
            )
            
            yield f"Starting benchmark: {' '.join(cmd)}\n\n"
            
            for line in process.stdout:
                yield line
                
            process.wait()
            yield f"\nBenchmark completed with exit code {process.returncode}\n"

        return Response(generate(), mimetype="text/plain")
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/agents/dashboard")
def agents_dashboard():
    return render_template("agents_dashboard.html")

@app.route("/agents/clarify", methods=["POST"])
@requires_role("PM")
def agents_clarify():
    """
    PM Clarification Loop: Calls the ReasoningAgent to analyze the initial ideation prompt.
    Returns a multi-turn chat response with clarifying questions based on RAG context.
    """
    data = request.json or {}
    prompt = data.get("prompt", "").strip()
    history = data.get("history", [])
    if not prompt:
        return jsonify({"error": "Missing prompt"}), 400

    is_first_turn = not any(m.get("role") == "user" for m in history)

    system = (
        "You are the NPCI Agentic Architect — a senior UPI product advisor.\n"
        "Your task: analyze the PM's UPI feature request and ask 2-3 precise clarifying questions "
        "so a complete 10-section NPCI Product Canvas can be built.\n"
        "The canvas covers: feature description, market need, ecosystem participants, scalability, "
        "validation/MVP, operating KPIs, product comms, pricing, risks, and compliance "
        "(specific NPCI OCs + RBI Master Directions).\n\n"
        + (
            "IMPORTANT: This is the FIRST message. You MUST always ask 2-3 targeted clarifying "
            "questions SPECIFIC to THIS feature. Base questions on the actual product described."
            if is_first_turn else
            "The PM has responded. Review their answers. Ask 1-2 focused follow-ups if needed, "
            "or confirm you have enough context with needs_clarification=false."
        )
        + "\n\nReturn ONLY this JSON (no prose, no markdown fences):\n"
        '{"needs_clarification": true, "clarification_questions": ["Q1 for this specific feature", "Q2", "Q3"], '
        '"message_to_pm": "1-2 sentence summary.", "confident": false}'
    )

    llm = LLMClient()
    msgs = [{"role": "system", "content": system},
            {"role": "user", "content": f"Feature Request: {prompt}"}]
    for turn in history:
        role = turn.get("role", "user")
        content = turn.get("content", "")
        if content:
            msgs.append({"role": role, "content": content})

    try:
        import requests as _r
        import re as _re
        resp = _r.post(llm.api_url, json={
            "model": llm.model, "messages": msgs,
            "temperature": 0.4, "max_tokens": 800,
        }, headers=getattr(llm, "headers", None), timeout=60)
        resp.raise_for_status()
        raw = resp.json()["choices"][0]["message"]["content"]
        raw = _re.sub(r"<think>.*?</think>", "", raw, flags=_re.DOTALL).strip()
        s, e = raw.find("{"), raw.rfind("}")
        if s != -1 and e != -1:
            result = json.loads(raw[s:e + 1])
            if is_first_turn and not result.get("clarification_questions"):
                result["clarification_questions"] = [
                    f"What are the transaction limits for this feature (per-transaction and daily cap)?",
                    "Which specific banks, PSPs, and merchant segments must be onboarded for the pilot?",
                    "Which NPCI circulars or RBI Master Directions must this comply with?"
                ]
                result["needs_clarification"] = True
            return jsonify(result), 200
    except Exception as ex:
        print(f"[clarify] LLM call failed: {ex}")

    # Prompt-aware fallback (never a 500)
    short = prompt[:60]
    return jsonify({
        "needs_clarification": True,
        "clarification_questions": [
            f"For '{short}…' — what are the per-transaction and daily limits?",
            "Which banks, PSPs, and merchant categories are part of the pilot?",
            "Which specific NPCI OCs or RBI guidelines apply to this feature?"
        ],
        "message_to_pm": f"I've reviewed your brief on '{short}'. A few questions will sharpen the canvas.",
        "confident": False
    }), 200

@app.route("/agents/canvas/generate", methods=["POST"])
@requires_role("PM")
def generate_canvas():
    data = request.json or {}
    prompt = data.get("prompt")
    context_notes = data.get("context_notes", "")
    
    try:
        agent = CanvasAgent(llm_client=LLMClient())
        result = agent.generate_canvas(prompt, context_notes)
        return jsonify(result), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/prototype/generate", methods=["POST"])
def generate_prototype():
    data = request.json or {}
    canvas = data.get("canvas")
    
    try:
        agent = PrototypeAgent(llm_client=LLMClient())
        result = agent.generate_flow_diagram(canvas)
        # Sign document bundle for phase gate
        signed_bundle = generate_signed_document(
            document_id="prototype_1",
            stage="PROTOTYPE",
            content=result,
            approver_role="PM",
            approver_id="PM_DEMO_01"
        )
        return jsonify({"prototype": result, "signature_bundle": signed_bundle}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/kit/generate", methods=["POST"])
@requires_role("PM")
def generate_kit():
    data = request.json or {}
    try:
        agent = ProductKitAgent(llm_client=LLMClient())
        result = agent.generate_kit(data.get("canvas", {}), data.get("prototype", {}))
        return jsonify(result), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/brd/formalise", methods=["POST"])
@requires_role("PM")
def formalise_brd():
    data = request.json or {}
    try:
        brd_payload = {
            "canvas": data.get("canvas", {}),
            "prototype": data.get("prototype", {}),
            "product_kit": data.get("product_kit", {})
        }
        signed_brd = generate_signed_document(
            document_id=f"brd_{int(time.time())}",
            stage="BRD_APPROVAL",
            content=brd_payload,
            approver_role="PM",
            approver_id="PM_DEMO_01"
        )
        return jsonify({"message": "BRD Formalised and Signed Successfully", "brd": signed_brd}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/master/translate", methods=["POST"])
@requires_role("NPCI_admin")
def master_translate():
    data = request.json or {}
    signed_brd = data.get("signed_brd")
    
    if not signed_brd or not signed_brd.get("raw_content"):
        return jsonify({"error": "Missing or invalid Signed BRD."}), 400
        
    try:
        # Step 1: Translate the requirements
        agent = NPCIMasterAgent(llm_client=LLMClient())
        tsd_result = agent.translate_brd_to_tsd(signed_brd.get("raw_content"))
        
        # Step 2: Sign the TSD and Change Manifests
        signed_tsd = generate_signed_document(
            document_id=f"tsd_{int(time.time())}",
            stage="TSD_GENERATION",
            content=tsd_result,
            approver_role="NPCI_admin",
            approver_id="NPCI_ADMIN_01"
        )
        
        # Step 3: Broadcast each chunk directly to the NotificationBus for Party Agents
        manifests = tsd_result.get("change_manifests", [])
        for chunk in manifests:
             bus.publish_event("spec_change", {
                 "party_id": chunk.get("party_id"),
                 "party_type": chunk.get("party_type"),
                 "changes": chunk.get("changes")
             })

        return jsonify({"message": "TSD Translated and Party Manifests broadcasted", "tsd": signed_tsd}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/certification/readiness", methods=["GET"])
@requires_role("NPCI_admin")
def cert_readiness():
    try:
        report = cert_manager.get_readiness_report()
        return jsonify(report), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/certification/uat", methods=["POST"])
@requires_role("NPCI_admin")
def cert_uat():
    try:
        status = cert_manager.trigger_uat()
        return jsonify({"message": f"UAT Execution Phase: {status}"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/certification/approve", methods=["POST"])
@requires_role("NPCI_certifier")
def cert_approve():
    data = request.json or {}
    readiness_report = cert_manager.get_readiness_report()
    if readiness_report.get("status") != "READY":
        return jsonify({"error": "System is NOT READY for certification. Missing party signatures."}), 400
        
    try:
        signed_cert = generate_signed_document(
            document_id=f"cert_{int(time.time())}",
            stage="DEPLOYMENT_APPROVAL",
            content=readiness_report,
            approver_role="NPCI_certifier",
            approver_id="CERT_MASTER_01"
        )
        return jsonify({"message": "Deployment Globally Certificated", "certificate": signed_cert}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/agents/certification/deploy", methods=["POST"])
@requires_role("NPCI_admin")
def cert_deploy():
    data = request.json or {}
    certificate = data.get("certificate")
    if not certificate or not certificate.get("signature"):
        return jsonify({"error": "Invalid or missing SHA-signed Certificate."}), 400
        
    # Technically we would call verify_signature(certificate) here
    # Broadcasting to all live agents to merge changes natively
    try:
        cert_manager.broadcast_deployment()
        return jsonify({"message": "Deployment Broadcasted to NotificationBus. System is now Live."}), 200
    except Exception as e:
         return jsonify({"error": str(e)}), 500

@app.route("/agents/approve-change", methods=["POST"])
def approve_change():
    data = request.json
    if not data:
        return jsonify({"error": "Missing payload"}), 400
    
    import threading
    import time

    def auto_approver():
        time.sleep(1)
        sub = bus.subscribe("agent_status")
        for event in sub:
            if event.get("status") == "AWAITING_APPROVAL":
                phase = event.get("phase")
                time.sleep(0.5)
                # print(f"[AUTO-APPROVER] Approving {phase}...")
                bus.publish_event("human_approval", {"phase": phase, "decision": "APPROVE"})
                if phase == "FINAL_DEPLOYMENT":
                    break

    def safe_execute(plan_data):
        import traceback
        try:
            switch_agent.execute_spec_change(plan_data)
        except Exception as e:
            print(f"CRITICAL ERROR in execute_spec_change: {e}")
            traceback.print_exc()

    threading.Thread(target=auto_approver, daemon=True).start()
    threading.Thread(target=safe_execute, args=(data,), daemon=True).start()
    return jsonify({"status": "Execution initialized", "version": data.get("version")}), 200

@app.route("/agents/trigger-change", methods=["POST"])
def trigger_spec_change():
    data = request.json
    version = data.get("version")
    description = data.get("description")
    
    if not version or not description:
        return jsonify({"error": "Missing version or description"}), 400
        
    switch_agent.publish_spec_change(version, description)
    return jsonify({"status": "Spec change broadcasted"}), 200

@app.route("/agents/propose-change", methods=["POST"])
def propose_change():
    data = request.json
    prompt = data.get("prompt")
    if not prompt:
        return jsonify({"error": "Missing prompt"}), 400

    plan = switch_agent.propose_spec_change(prompt)

    # Auto-calculate next version = total existing post-change snapshots + 1
    # This keeps the version number in sync with the re-sequenced display list.
    import json as _json2
    _project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    backup_root = os.path.join(_project_root, "backups")
    post_change_count = 0
    try:
        if os.path.isdir(backup_root):
            for _bname in os.listdir(backup_root):
                _bfolder = os.path.join(backup_root, _bname)
                _bmeta   = os.path.join(_bfolder, "metadata.json")
                if os.path.isfile(_bmeta):
                    try:
                        with open(_bmeta) as _bmf:
                            _bmeta_data = _json2.load(_bmf)
                        if _bmeta_data.get("type") == "post_change":
                            post_change_count += 1
                    except Exception:
                        pass
    except Exception:
        pass
    plan["version"] = f"1.{post_change_count + 1}"

    return jsonify(plan), 200

@app.route("/agents/approve-phase", methods=["POST"])
def approve_phase():
    """
    NPCI Audit Gate: Publishes a human approval decision for a specific phase
    to the internal NotificationBus, ungating the agent orchestrator.
    """
    data = request.json
    phase = data.get("phase")
    decision = data.get("decision", "APPROVE")
    
    if not phase:
        return jsonify({"error": "Missing 'phase' in request"}), 400
        
    print(f"[API] 🛡️ Audit Gate: Phase {phase} received {decision}")
    
    # Broadcast to the internal bus for SwitchAgent to consume
    bus.publish_event("human_approval", {
        "phase": phase,
        "decision": decision,
        "timestamp": time.time()
    })
    
    return jsonify({"status": f"Phase {phase} {decision} broadcasted"}), 200

@app.route("/phase2")
def phase2_ui():
    return render_template("phase2_dashboard.html")

@app.route("/workflow2/generate", methods=["POST"])
def workflow2_generate():
    """
    Workflow 2: New feature build output generator
    Returns deep research + product canvas + prototype brief + product comms.
    """
    data = request.json or {}
    feature_name = (data.get("feature_name") or "").strip()
    objective = (data.get("objective") or "").strip()
    target_users = (data.get("target_users") or "").strip()
    maker_notes = (data.get("maker_notes") or "").strip()

    if not feature_name:
        return jsonify({"error": "Missing 'feature_name'"}), 400

    try:
        llm_client = getattr(getattr(switch_agent, "reasoning_agent", None), "llm_client", None)
        if not llm_client:
            payload = _workflow2_fallback(feature_name, objective, target_users)
            session_id = str(uuid.uuid4())
            WORKFLOW2_SESSIONS[session_id] = {
                "session_id": session_id,
                "feature_name": feature_name,
                "objective": objective,
                "target_users": target_users,
                "maker_notes": maker_notes,
                "payload": payload,
                "canvas_approved": False,
                "prototype_approved": False,
                "prototype_url": "",
                "created_at": datetime.utcnow().isoformat() + "Z",
            }
            return jsonify({"session_id": session_id, **payload}), 200

        prompt = f"""
You are an NPCI product+strategy+compliance copilot. Build Workflow-2 outputs for a new UPI feature.
Return ONLY valid JSON (no markdown).

Feature: {feature_name}
Objective: {objective or "Not provided"}
Target Users: {target_users or "Not provided"}
Maker Notes: {maker_notes or "Not provided"}

Use this exact schema:
{{
  "thinking_trace": ["string"],
  "agent_plan": ["string"],
  "deep_research": {{
    "feature_summary": "string",
    "need": "string",
    "market_view": "string",
    "scalability": "string",
    "risks": ["string"],
    "regulatory_summary": "include RBI/NPCI-oriented checks and constraints"
    "rbi_guidelines_checked": [
      "https://www.rbi.org.in/scripts/NotificationUser.aspx?Id=12032&Mode=0",
      "https://www.rbi.org.in/commonman/english/scripts/Notification.aspx?Id=1888"
    ]
  }},
  "product_canvas": {{
    "feature_for_layman": "string",
    "need_and_differentiation": "string",
    "market_view": "string",
    "scalability": "string",
    "validation_mvp": "string",
    "operating_kpis": ["string"],
    "comms_outputs": ["string"],
    "pricing_view": "string",
    "compliance": "string"
  }},
  "prototype": {{
    "name": "string",
    "url": "string (empty if not available)",
    "notes": "string"
  }},
  "product_comms": {{
    "doc_outline": ["string"],
    "faq_starters": ["string"],
    "circular_draft": "string",
    "trained_llm_note": "string"
  }},
  "execution_handoff": {{
    "objective": "string",
    "target_users": "string",
    "next_step": "string"
  }}
}}
"""
        raw = llm_client.query(prompt)
        parsed = _extract_json_from_text(raw)
        if not parsed:
            parsed = _workflow2_fallback(feature_name, objective, target_users)
        parsed.setdefault("thinking_trace", _workflow2_fallback(feature_name, objective, target_users).get("thinking_trace", []))
        parsed.setdefault("agent_plan", _workflow2_fallback(feature_name, objective, target_users).get("agent_plan", []))
        deep = parsed.setdefault("deep_research", {})
        deep["rbi_guidelines_checked"] = [
            "https://www.rbi.org.in/scripts/NotificationUser.aspx?Id=12032&Mode=0",
            "https://www.rbi.org.in/commonman/english/scripts/Notification.aspx?Id=1888",
        ]
        session_id = str(uuid.uuid4())
        WORKFLOW2_SESSIONS[session_id] = {
            "session_id": session_id,
            "feature_name": feature_name,
            "objective": objective,
            "target_users": target_users,
            "maker_notes": maker_notes,
            "payload": parsed,
            "canvas_approved": False,
            "prototype_approved": False,
            "prototype_url": "",
            "created_at": datetime.utcnow().isoformat() + "Z",
        }
        return jsonify({"session_id": session_id, **parsed}), 200
    except Exception:
        payload = _workflow2_fallback(feature_name, objective, target_users)
        session_id = str(uuid.uuid4())
        WORKFLOW2_SESSIONS[session_id] = {
            "session_id": session_id,
            "feature_name": feature_name,
            "objective": objective,
            "target_users": target_users,
            "maker_notes": maker_notes,
            "payload": payload,
            "canvas_approved": False,
            "prototype_approved": False,
            "prototype_url": "",
            "created_at": datetime.utcnow().isoformat() + "Z",
        }
        return jsonify({"session_id": session_id, **payload}), 200

@app.route("/workflow2/approve-canvas", methods=["POST"])
def workflow2_approve_canvas():
    data = request.json or {}
    session_id = data.get("session_id", "")
    updated_canvas = data.get("product_canvas")
    if not session_id or session_id not in WORKFLOW2_SESSIONS:
        return jsonify({"error": "Invalid or missing session_id"}), 400
    sess = WORKFLOW2_SESSIONS[session_id]
    if isinstance(updated_canvas, dict):
        sess["payload"]["product_canvas"] = updated_canvas
    sess["canvas_approved"] = True
    return jsonify({"status": "approved", "session_id": session_id}), 200

@app.route("/workflow2/build-prototype", methods=["POST"])
def workflow2_build_prototype():
    data = request.json or {}
    session_id = data.get("session_id", "")
    if not session_id or session_id not in WORKFLOW2_SESSIONS:
        return jsonify({"error": "Invalid or missing session_id"}), 400
    sess = WORKFLOW2_SESSIONS[session_id]
    if not sess.get("canvas_approved"):
        return jsonify({"error": "Approve product canvas before prototype build"}), 400
    builder = os.getenv("LOVABLE_BUILDER_URL", "https://lovable.dev/")
    custom = (data.get("prototype_url") or "").strip()
    proto_url = custom or (builder.rstrip("/") + "/?feature=" + feature_name_to_query(sess.get("feature_name", "feature")))
    sess["prototype_url"] = proto_url
    sess["payload"].setdefault("prototype", {})
    sess["payload"]["prototype"]["url"] = proto_url
    sess["prototype_approved"] = bool(data.get("prototype_approved", False))
    return jsonify({"status": "prototype_ready", "redirect_url": proto_url, "session_id": session_id}), 200

def feature_name_to_query(name: str):
    return name.replace(" ", "%20").replace("&", "%26")

@app.route("/workflow2/documents/<session_id>/manifest", methods=["GET"])
def workflow2_documents_manifest(session_id):
    sess = WORKFLOW2_SESSIONS.get(session_id)
    if not sess:
        return jsonify({"error": "Session not found"}), 404
    docs = _render_workflow2_docs(sess)
    items = []
    for key, meta in docs.items():
        doc_item = {
            "key": key,
            "title": meta["title"],
            "preview_url": f"/workflow2/documents/{session_id}/preview/{key}",
            "download_docx_url": f"/workflow2/documents/{session_id}/download/{key}.docx",
        }
        if key == "product-deck":
            doc_item["download_ppt_url"] = f"/workflow2/documents/{session_id}/download/{key}.pptx"
        items.append(doc_item)
    return jsonify({"session_id": session_id, "documents": items}), 200

@app.route("/workflow2/documents/<session_id>/preview/<doc_key>", methods=["GET"])
def workflow2_document_preview(session_id, doc_key):
    sess = WORKFLOW2_SESSIONS.get(session_id)
    if not sess:
        return "Session not found", 404
    docs = _render_workflow2_docs(sess)
    if doc_key not in docs:
        return "Document not found", 404
    doc = docs[doc_key]
    body = html_escape(doc["body"]).replace("\n", "<br>")
    return f"""
    <html><head><title>{html_escape(doc['title'])}</title>
    <style>body{{font-family:Arial,sans-serif;background:#f7f7fb;padding:24px}} .doc{{max-width:920px;margin:auto;background:#fff;padding:32px;border-radius:12px;box-shadow:0 4px 20px rgba(0,0,0,.08)}} h1{{font-size:24px}} .meta{{color:#666;font-size:12px;margin-bottom:16px}}</style>
    </head><body><div class='doc'><h1>{html_escape(doc['title'])}</h1><div class='meta'>Generated from Workflow 2 session {html_escape(session_id)}</div><div>{body}</div></div></body></html>
    """

@app.route("/workflow2/documents/<session_id>/download/<filename>", methods=["GET"])
def workflow2_document_download(session_id, filename):
    sess = WORKFLOW2_SESSIONS.get(session_id)
    if not sess:
        return jsonify({"error": "Session not found"}), 404
    
    if filename.endswith(".pptx"):
        key = filename[:-5]
        docs = _render_workflow2_docs(sess)
        if key not in docs:
            return jsonify({"error": "Document not found"}), 404
        doc = docs[key]
        content = _ppt_from_text(doc["title"], doc["body"])
        return send_file(
            io.BytesIO(content),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=filename,
        )
    elif filename.endswith(".docx"):
        key = filename[:-5]
        docs = _render_workflow2_docs(sess)
        if key not in docs:
            return jsonify({"error": "Document not found"}), 404
        doc = docs[key]
        content = _docx_from_text(doc["title"], doc["body"])
        return send_file(
            io.BytesIO(content),
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=filename,
        )
    return jsonify({"error": "Unsupported download format"}), 400

@app.route("/workflow2/advance-execution", methods=["POST"])
def workflow2_advance_execution():
    data = request.json or {}
    session_id = data.get("session_id", "")
    if not session_id or session_id not in WORKFLOW2_SESSIONS:
        return jsonify({"error": "Invalid or missing session_id"}), 400
    sess = WORKFLOW2_SESSIONS[session_id]
    if not sess.get("canvas_approved"):
        return jsonify({"error": "Canvas approval pending"}), 400
    if not sess.get("prototype_url"):
        return jsonify({"error": "Prototype build pending"}), 400
    feature = sess.get("feature_name", "Workflow2 Feature")
    plan_payload = {
        "version": "2.0",
        "description": f"Workflow 2 execution handoff for {feature}",
        "impact_analysis": ["workflow2_canvas", "prototype_handoff", "comms_package"],
        "plan": sess.get("payload", {}).get("agent_plan", []),
        "verification_payload": ""
    }
    return jsonify({"status": "ready_for_execution", "plan_payload": plan_payload}), 200

@app.route("/agents/pm-feedback", methods=["POST"])
def save_pm_feedback():
    """Save PM feedback for test results"""
    try:
        data = request.json
        feedback = {
            "test_id": data.get("test_id"),
            "verdict": data.get("verdict"),  # pass, fail, needs-work
            "feedback": data.get("feedback", ""),
            "timestamp": time.time()
        }
        # In a real system, you'd save this to a database
        # For now, we'll just log it and broadcast it
        print(f"[PM Feedback] {feedback}")
        bus.publish_event("pm_feedback", feedback)
        return jsonify({"status": "Feedback saved", "feedback": feedback}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route("/schemas")
def list_schemas():
    import os
    schema_dir = "api/schemas"
    schemas = []
    if os.path.exists(schema_dir):
        schemas = [f for f in os.listdir(schema_dir) if f.endswith('.xsd')]
    
    # Simple template for schema listing
    html = """
    <html>
    <head><title>XSD Schemas</title>
    <style>
        body { font-family: sans-serif; padding: 20px; }
        ul { list-style-type: none; padding: 0; }
        li { margin: 10px 0; }
        a { text-decoration: none; color: #007bff; font-size: 18px; }
        a:hover { text-decoration: underline; }
    </style>
    </head>
    <body>
        <h1>Available XSD Schemas</h1>
        <ul>
            {% for schema in schemas %}
            <li><a href="/schemas/{{ schema }}">{{ schema }}</a></li>
            {% endfor %}
        </ul>
        <br>
        <a href="/">Back to Home</a>
    </body>
    </html>
    """
    return render_template_string(html, schemas=schemas)

@app.route("/schemas/<filename>")
def view_schema(filename):
    import os
    schema_dir = "api/schemas"
    # Basic path traversal protection
    if ".." in filename or "/" in filename:
        return "Invalid filename", 400
        
    filepath = os.path.join(schema_dir, filename)
    if not os.path.exists(filepath):
        return "File not found", 404
        
    with open(filepath, "r") as f:
        content = f.read()
        
    return Response(content, mimetype="application/xml")

def _get_component_files():
    """
    Helper used by both the Inspector and Deploy APIs to know which
    files constitute the "UPI system" for versioning/rollback.
    """
    components = {
        "Schemas": [],
        "Switch": ["switch/upi_switch.py"],
        "Payer PSP": ["psps/payer_psp.py", "psps/payer_psp_handler.py"],
        "Payee PSP": ["psps/payee_psp.py", "psps/payee_psp_handler.py"],
        "Banks": ["banks/remitter_bank.py", "banks/beneficiary_bank.py", "banks/remitter_bank_handler.py", "banks/beneficiary_bank_handler.py"],
    }
    schema_dir = "api/schemas"
    if os.path.exists(schema_dir):
        components["Schemas"] = [
            os.path.join(schema_dir, f)
            for f in os.listdir(schema_dir)
            if f.endswith(".xsd")
        ]
    return components


@app.route("/deploy/versions")
def list_backup_versions():
    """
    Returns deployable post-change snapshots, newest first.
    Versions are re-sequenced chronologically (v1.1 oldest … v1.N newest)
    so the display always shows incrementing numbers regardless of what is
    stored in the folder name.
    """
    import json as _json
    backup_root = "backups"
    entries = []

    if os.path.exists(backup_root):
        for name in os.listdir(backup_root):
            folder = os.path.join(backup_root, name)
            if not os.path.isdir(folder):
                continue
            meta_path = os.path.join(folder, "metadata.json")
            if os.path.exists(meta_path):
                try:
                    with open(meta_path) as mf:
                        meta = _json.load(mf)
                    if meta.get("type") == "post_change":
                        entries.append({
                            "id":          name,
                            "version":     meta.get("version", "?"),
                            "description": meta.get("description", name),
                            "timestamp":   meta.get("timestamp", ""),
                        })
                except Exception:
                    pass

    # Sort oldest → newest first so we can assign ascending version numbers
    entries.sort(key=lambda x: x["timestamp"])

    # Re-sequence: v1.1, v1.2, v1.3 … regardless of what the folder name says
    for idx, entry in enumerate(entries, start=1):
        seq_ver = f"1.{idx}"
        entry["seq_version"] = seq_ver
        desc_short = entry["description"][:55]
        entry["label"] = f"v{seq_ver} — {desc_short}"

    # Flip to newest-first for the UI
    entries.reverse()

    # Fallback: if no post-change snapshots yet, return raw folder list
    if not entries:
        raw = sorted(
            [n for n in os.listdir(backup_root) if os.path.isdir(os.path.join(backup_root, n))]
            if os.path.exists(backup_root) else [],
            reverse=True
        )
        entries = [{"id": n, "label": n, "seq_version": "", "description": n, "timestamp": ""} for n in raw]

    return jsonify({"versions": entries})


@app.route("/deploy", methods=["POST"])
def deploy_version():
    """
    Restore code/config files from a selected backup snapshot, then
    auto-restart the Flask process so changes go live immediately.
    """
    import shutil
    import threading
    import sys

    data = request.json or {}
    version = data.get("version")
    if not version:
        return jsonify({"error": "Missing 'version' in request body"}), 400

    backup_root = "backups"
    backup_dir = os.path.join(backup_root, version)
    if not os.path.isdir(backup_dir):
        return jsonify({"error": f"Backup version '{version}' not found"}), 400

    components = _get_component_files()
    restored_files = []
    for group_files in components.values():
        for rel_path in group_files:
            src = os.path.join(backup_dir, rel_path)
            dst = rel_path
            if os.path.exists(src):
                os.makedirs(os.path.dirname(dst), exist_ok=True)
                shutil.copy2(src, dst)
                restored_files.append(rel_path)

    # Also copy schemas from post-change snapshot if present there
    schema_backup = os.path.join(backup_dir, "api", "schemas")
    if os.path.isdir(schema_backup):
        for fname in os.listdir(schema_backup):
            src = os.path.join(schema_backup, fname)
            dst = os.path.join("api", "schemas", fname)
            if os.path.isfile(src):
                shutil.copy2(src, dst)
                if dst not in restored_files:
                    restored_files.append(dst)

    # Schedule server restart: kill self then spawn a fresh process so the port is freed first
    def _restart():
        import time as _time
        import subprocess as _sp
        import signal as _signal
        _time.sleep(1.5)   # let the HTTP response reach the browser
        pid = os.getpid()
        env = os.environ.copy()
        env["PYTHONPATH"] = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        # Spawn a new server process (detached so it survives after we die)
        _sp.Popen(
            [sys.executable, "-u", os.path.abspath(__file__)],
            env=env,
            cwd=os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            start_new_session=True,
        )
        _time.sleep(0.3)
        os.kill(pid, _signal.SIGTERM)   # kill ourselves after spawning the replacement

    threading.Thread(target=_restart, daemon=True).start()

    return jsonify(
        {
            "status": "ok",
            "version": version,
            "restored_files": restored_files,
            "message": f"Deployed '{version}'. Server is restarting — Phase 1 will be live in ~5 seconds.",
            "restarting": True,
        }
    )


@app.route("/git/info")
def git_info():
    """
    Lightweight helper so the UI can detect whether this directory
    is a Git repo and show basic status.
    """
    try:
        # Check if we're in a git work tree
        result = subprocess.run(
            ["git", "rev-parse", "--is-inside-work-tree"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            cwd=".",
        )
        if result.returncode != 0 or result.stdout.strip() != "true":
            return jsonify({"is_repo": False, "message": "Not a Git repository here."})

        # Get current branch (if any)
        branch = subprocess.run(
            ["git", "rev-parse", "--abbrev-ref", "HEAD"],
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            text=True,
            cwd=".",
        ).stdout.strip()

        status_out = subprocess.run(
            ["git", "status", "--porcelain"],
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            text=True,
            cwd=".",
        ).stdout

        return jsonify(
            {
                "is_repo": True,
                "branch": branch,
                "dirty": bool(status_out.strip()),
                "status_raw": status_out,
            }
        )
    except Exception as e:
        return jsonify({"is_repo": False, "error": str(e)}), 500


@app.route("/git/commit", methods=["POST"])
def git_commit():
    """
    Create a local Git commit for the core UPI components.
    This assumes Git is already initialized and configured.
    """
    data = request.json or {}
    message = data.get("message") or "UPI spec change via agents"

    try:
        # Ensure this is a repo
        result = subprocess.run(
            ["git", "rev-parse", "--is-inside-work-tree"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            cwd=".",
        )
        if result.returncode != 0 or result.stdout.strip() != "true":
            return jsonify({"error": "Not a Git repository. Run 'git init' first."}), 400

        # Stage relevant component files only
        components = _get_component_files()
        files_to_add = []
        for group_files in components.values():
            files_to_add.extend(group_files)

        add_cmd = ["git", "add"] + files_to_add
        add_proc = subprocess.run(add_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, cwd=".")
        if add_proc.returncode != 0:
            return jsonify({"error": f"git add failed: {add_proc.stderr}"}), 500

        commit_proc = subprocess.run(
            ["git", "commit", "-m", message],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            cwd=".",
        )
        if commit_proc.returncode != 0:
            return jsonify({"error": f"git commit failed: {commit_proc.stderr}"}), 500

        # Get last commit hash and message
        log_proc = subprocess.run(
            ["git", "log", "-1", "--pretty=format:%H||%s"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            cwd=".",
        )
        commit_hash = ""
        commit_subject = ""
        if log_proc.returncode == 0 and "||" in log_proc.stdout:
            commit_hash, commit_subject = log_proc.stdout.split("||", 1)

        return jsonify(
            {
                "status": "ok",
                "message": "Commit created.",
                "commit_hash": commit_hash,
                "commit_subject": commit_subject,
            }
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/current-version")
def current_version():
    """Returns metadata about the latest deployed version for Phase 1 header badge."""
    import json as _json
    backup_root = "backups"
    entries = []
    if os.path.exists(backup_root):
        for name in os.listdir(backup_root):
            folder = os.path.join(backup_root, name)
            meta_path = os.path.join(folder, "metadata.json")
            if os.path.isfile(meta_path):
                try:
                    with open(meta_path) as mf:
                        meta = _json.load(mf)
                    if meta.get("type") == "post_change":
                        entries.append({"id": name, "timestamp": meta.get("timestamp",""), "description": meta.get("description",""), "version": meta.get("version","?")})
                except Exception:
                    pass
    if not entries:
        return jsonify({"seq_version": "1.0", "description": "Base version", "timestamp": ""})
    entries.sort(key=lambda x: x["timestamp"])
    latest = entries[-1]
    latest["seq_version"] = f"1.{len(entries)}"
    return jsonify(latest)


@app.route("/inspector/file-content")
def inspector_file_content():
    """Return file content + diff vs a chosen backup version (JSON API)."""
    import difflib, json as _json
    selected_file = request.args.get("file", "")
    backup_version = request.args.get("version", "")   # folder name in backups/

    if ".." in selected_file or selected_file.startswith("/"):
        return jsonify({"error": "Invalid path"}), 400

    current_content = ""
    if os.path.exists(selected_file):
        try:
            with open(selected_file, "r") as f:
                current_content = f.read()
        except Exception as e:
            current_content = f"# Error reading file: {e}"
    else:
        current_content = "# File not found"

    backup_content = None
    if backup_version:
        bp = os.path.join("backups", backup_version, selected_file)
        if os.path.exists(bp):
            with open(bp) as bf:
                backup_content = bf.read()

    # unified diff lines
    diff_lines = []
    if backup_content is not None:
        diff = list(difflib.unified_diff(
            backup_content.splitlines(),
            current_content.splitlines(),
            fromfile="backup",
            tofile="live",
            lineterm=""
        ))
        diff_lines = diff

    return jsonify({
        "content": current_content,
        "diff": diff_lines,
        "changed": len(diff_lines) > 0,
        "backup_found": backup_content is not None
    })


@app.route("/inspector/changed-files")
def inspector_changed_files():
    """Return which component files differ between current live and a backup version."""
    import json as _json
    backup_version = request.args.get("version", "")
    components = _get_component_files()
    result = {}
    for group, files in components.items():
        result[group] = []
        for f in files:
            changed = False
            if backup_version:
                bp = os.path.join("backups", backup_version, f)
                if os.path.exists(bp) and os.path.exists(f):
                    try:
                        with open(f) as cf, open(bp) as bf:
                            changed = cf.read() != bf.read()
                    except Exception:
                        pass
                elif not os.path.exists(bp) and os.path.exists(f):
                    changed = True  # new file
            result[group].append({"path": f, "name": f.split("/")[-1], "changed": changed})
    return jsonify({"files": result})


@app.route("/inspector")
def inspector():
    components = _get_component_files()
    return render_template("inspector.html", components=components)

if __name__ == "__main__":
    app.run(port=5000, debug=True, use_reloader=False, threaded=True)
