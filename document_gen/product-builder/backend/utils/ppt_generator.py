import io
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt(title: str, content: str, path: str):
    """
    Generate a 5-slide formal PPT deck from markdown-ish content.
    """
    prs = Presentation()
    
    # Split content by ## Slide headers
    sections = []
    if "## Slide" in content:
        parts = content.split("## Slide")
        for p in parts[1:]:
            sections.append(p.strip())
    else:
        sections = [content]

    slide_titles = [
        "Product Overview",
        "Market Dynamics",
        "Product Perspective",
        "Strategic Roadmap",
        "Governance & Compliance"
    ]
    
    for i in range(5):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_titles[i]
        
        # Body
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Content
        txt = sections[i] if i < len(sections) else "Details pending refinement."
        # Strip the first line if it's the slide number title
        lines = txt.split("\n", 1)
        if len(lines) > 1 and ":" in lines[0]:
            txt = lines[1].strip()
            
        tf.text = txt

    prs.save(path)
